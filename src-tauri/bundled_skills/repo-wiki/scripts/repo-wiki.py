#!/usr/bin/env python3
"""repo-wiki.py — Unified repository wiki system.

Bi-directional linking + binary conversion + full indexing.

Usage:
    # Wiki (metadata, links, backlinks)
    python3 repo-wiki.py wiki init
    python3 repo-wiki.py wiki add [--all] [--file PATH]
    python3 repo-wiki.py wiki link <from-file> <to-slug> [--section HEADING]
    python3 repo-wiki.py wiki backlinks --build [--file PATH]
    python3 repo-wiki.py wiki links --check [--fail-on-broken]
    python3 repo-wiki.py wiki migrate --to STRUCTURED
    python3 repo-wiki.py wiki status

    # Binary conversion
    python3 repo-wiki.py convert [--root PATH] [--formats EXTENSIONS] [--out DIR] [--overwrite]

    # Index generation
    python3 repo-wiki.py index [--root PATH] [--out PATH] [--dry-run] [--ignore-dirs DIRS]

    # Unified runner
    python3 repo-wiki.py run [--root PATH] [--skip-convert] [--skip-index] [--dry-run]

    # Dependencies
    python3 repo-wiki.py install-deps [--fmt EXTENSIONS]
"""

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SKILL_DIR = Path(__file__).resolve().parent
REPO_ROOT = SKILL_DIR.parent.parent.parent.parent
DEFAULT_DOCS_DIR = REPO_ROOT / "docs"

FRONTMATTER_RE = re.compile(r"^---\n(.*?)\n---\n", re.DOTALL)
SLUG_RE = re.compile(r"\[\[([a-z0-9][a-z0-9-]*)(?:#([a-z0-9_-]+))?\]\]", re.IGNORECASE)
RELATIVE_LINK_RE = re.compile(r"\[([^\]]+)\]\(\.?/?(?:docs/)?([^\)#\s]+\.md)(?:#([^\)#\s]+))?\)")
HEADING_RE = re.compile(r"^(#{1,6})\s+(.+)$", re.MULTILINE)

DEFAULT_STATUS = "draft"
DEFAULT_CATEGORY = "general"
STRUCTURED_DIRS = {
    "core": ["architecture", "guides", "contributing", "api", "features"],
    "reference": ["mcp", "llm-services", "3rd_party"],
    "drafts": ["improvements", "refactoring"],
    "archive": ["analysis", "fixes", "sprints", "tauri", "testing"],
}

IGNORE_DIRS = {".git", "__pycache__", "node_modules", ".github", "venv"}
BINARY_FORMATS = {"pdf": "pymupdf", "docx": "python-docx", "pptx": "python-pptx", "xlsx": "openpyxl"}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def slug_from_path(rel_path: str) -> str:
    name = Path(rel_path).stem
    name = re.sub(r"[_\s]+", "-", name).lower()
    name = re.sub(r"[^a-z0-9-]", "", name)
    name = re.sub(r"-{2,}", "-", name)
    return name.strip("-")


def heading_to_slug(heading: str) -> str:
    slug = heading.lower()
    slug = re.sub(r"[^a-z0-9\s-]", "", slug)
    slug = re.sub(r"\s+", "-", slug)
    return slug.strip("-")


def extract_frontmatter(content: str) -> tuple[dict[str, Any], str]:
    match = FRONTMATTER_RE.match(content)
    if not match:
        return {}, content
    fm_text = match.group(1)
    body = content[match.end():]
    metadata: dict[str, Any] = {}
    for line in fm_text.strip().split("\n"):
        line = line.strip()
        if ":" not in line:
            continue
        key, _, value = line.partition(":")
        key = key.strip()
        value = value.strip().strip('"').strip("'")
        if key == "tags":
            value = re.findall(r'"([^"]*)"', value) or re.findall(r"'([^']*)'", value) or [v.strip() for v in value.split(",") if v.strip()]
        metadata[key] = value
    return metadata, body


def write_frontmatter(path: Path, metadata: dict[str, Any], body: str) -> None:
    fm_lines = ["---"]
    for key, value in metadata.items():
        if isinstance(value, list):
            fm_lines.append(f"{key}: [{', '.join(f'\"{v}\"' for v in value)}]")
        else:
            fm_lines.append(f"{key}: {value}")
    fm_lines.append("---")
    path.write_text("\n".join(fm_lines) + "\n\n" + body.lstrip("\n"), encoding="utf-8")


def read_md(path: Path) -> tuple[dict[str, Any], str]:
    content = path.read_text(encoding="utf-8")
    return extract_frontmatter(content)


def find_md_files(directory: Path) -> list[Path]:
    files = []
    for root, dirs, filenames in os.walk(directory):
        root_path = Path(root)
        if any(d in root_path.parts for d in ("_meta",)):
            continue
        for fname in filenames:
            if fname.endswith(".md") and fname != "_index.md":
                files.append(Path(root) / fname)
    return sorted(files)


def find_binary_files(directory: Path, formats: list[str] = list(BINARY_FORMATS.keys())) -> list[Path]:
    exts = {f".{fmt}" for fmt in formats}
    files = []
    for root, dirs, filenames in os.walk(directory):
        root_path = Path(root)
        if any(d in root_path.parts for d in IGNORE_DIRS):
            continue
        for fname in filenames:
            if Path(fname).suffix.lower() in exts:
                files.append(Path(root) / fname)
    return sorted(files)


def ensure_dirs(*dirs: Path) -> None:
    for d in dirs:
        d.mkdir(parents=True, exist_ok=True)


def load_catalog(catalog_path: Path) -> dict[str, Any]:
    if catalog_path.exists():
        return json.loads(catalog_path.read_text(encoding="utf-8"))
    return {"version": "2.0", "generated": "", "docs": []}


def save_catalog(catalog: dict[str, Any], catalog_path: Path) -> None:
    catalog["generated"] = datetime.now(timezone.utc).isoformat()
    catalog_path.parent.mkdir(parents=True, exist_ok=True)
    catalog_path.write_text(json.dumps(catalog, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def get_doc_entry(slug: str, catalog: dict) -> Optional[dict]:
    for doc in catalog.get("docs", []):
        if doc.get("slug") == slug:
            return doc
    return None


# ===========================================================================
# Wiki Commands
# ===========================================================================

def cmd_wiki_init(root: Path, meta_dir: Path, backlinks_dir: Path, catalog_path: Path, index_path: Path, rc_path: Path) -> int:
    print("Initializing repo-wiki structure...")
    ensure_dirs(meta_dir, backlinks_dir)

    rc = {"version": "2.0", "repo_root": str(root), "docs_dir": str(root.relative_to(root.parent)),
          "meta_dir": "_meta", "default_status": DEFAULT_STATUS, "default_category": DEFAULT_CATEGORY}
    rc_path.write_text(json.dumps(rc, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    print(f"  Created .repo-wikirc")

    save_catalog({"version": "2.0", "generated": "", "docs": []}, catalog_path)
    print(f"  Created _meta/catalog.json")

    index_content = f"""# Repository Wiki Index

> Auto-generated by `repo-wiki.py init`. Edit to curate.

## Quick Links

| Category | Description |
|----------|-------------|
| **Architecture** | System design, data flow, security |
| **Guides** | Step-by-step tutorials |
| **API Reference** | Tauri commands, HTTP API |
| **Contributing** | Coding standards, release process |

## Recent Updates

<!-- Manual updates -->
"""
    index_path.write_text(index_content, encoding="utf-8")
    print(f"  Created _index.md")
    print("Done. Run `repo-wiki.py wiki add --all` to index all documents.")
    return 0


def cmd_wiki_add(root: Path, meta_dir: Path, catalog_path: Path, all_files: bool, file_path: Optional[str]) -> int:
    files_to_add: list[Path] = []
    if all_files:
        docs_dir = root / "docs"
        if docs_dir.is_dir():
            files_to_add = find_md_files(docs_dir)
        else:
            files_to_add = find_md_files(root)
    elif file_path:
        fpath = root / file_path
        if not fpath.exists():
            print(f"Error: File not found: {file_path}", file=sys.stderr)
            return 1
        files_to_add = [fpath]
    if not files_to_add:
        print("No files to add. Use --all or --file PATH.")
        return 0

    catalog = load_catalog(catalog_path)
    updated = created = 0

    for fpath in files_to_add:
        rel_path = str(fpath.relative_to(root))
        slug = slug_from_path(rel_path)
        metadata, body = read_md(fpath)

        if "slug" not in metadata:
            metadata["slug"] = slug
        if "title" not in metadata:
            first_heading = HEADING_RE.search(body)
            metadata["title"] = first_heading.group(2).strip() if first_heading else Path(rel_path).stem.replace("-", " ").title()
        if "status" not in metadata:
            metadata["status"] = DEFAULT_STATUS
        if "category" not in metadata:
            for cat, dirs in STRUCTURED_DIRS.items():
                if any(d in rel_path for d in dirs):
                    metadata["category"] = cat
                    break
            else:
                metadata["category"] = DEFAULT_CATEGORY
        if "tags" not in metadata:
            metadata["tags"] = []
        if "created" not in metadata:
            metadata["created"] = fpath.stat().st_mtime
        if "updated" not in metadata:
            metadata["updated"] = datetime.now(timezone.utc).strftime("%Y-%m-%d")

        old_content = fpath.read_text(encoding="utf-8")
        write_frontmatter(fpath, metadata, body)
        if old_content != fpath.read_text(encoding="utf-8"):
            if old_content.strip().startswith("---"):
                updated += 1
            else:
                created += 1

        entry = {"slug": metadata["slug"], "path": rel_path, "title": metadata["title"],
                 "status": metadata["status"], "category": metadata["category"],
                 "tags": metadata["tags"] if isinstance(metadata["tags"], list) else [],
                 "references": [], "referencedBy": [], "size": len(body), "updated": metadata["updated"]}
        catalog["docs"] = [d for d in catalog["docs"] if d.get("slug") != slug]
        catalog["docs"].append(entry)

    save_catalog(catalog, catalog_path)
    print(f"Added {created} new, updated {updated} existing. Catalog: {len(catalog['docs'])} total.")
    return 0


def cmd_wiki_link(root: Path, from_file: str, to_slug: str, section: Optional[str]) -> int:
    fpath = root / from_file
    if not fpath.exists():
        print(f"Error: File not found: {from_file}", file=sys.stderr)
        return 1

    target_ref = f"[[{to_slug}]]"
    if section:
        target_ref = f"[[{to_slug}#{section}]]"

    content = fpath.read_text(encoding="utf-8")
    catalog = load_catalog(root / "_meta" / "catalog.json")

    def replace_link(m):
        link_path = m.group(2)
        link_section = m.group(3)
        clean = link_path.replace("docs/", "")
        candidate = slug_from_path(clean)
        entry = get_doc_entry(candidate, catalog)
        if entry:
            ref = f"[[{entry['slug']}]]"
            if link_section or section:
                anchor = link_section or heading_to_slug(m.group(1))
                ref = f"[[{entry['slug']}#{anchor}]]"
            return f"{m.group(1)}{ref}"
        return m.group(0)

    new_content = RELATIVE_LINK_RE.sub(replace_link, content)
    if new_content != content:
        fpath.write_text(new_content, encoding="utf-8")
        print(f"Updated links in {from_file}")
    else:
        print(f"No matching links in {from_file}")
    return 0


def cmd_wiki_backlinks(root: Path, meta_dir: Path, catalog_path: Path, file_path: Optional[str]) -> int:
    catalog = load_catalog(catalog_path)
    docs_by_slug = {d["slug"]: d for d in catalog["docs"]}
    all_md = find_md_files(root)
    references: dict[str, list[str]] = {}

    for md_file in all_md:
        content = md_file.read_text(encoding="utf-8")
        for slug_match, section_match in SLUG_RE.findall(content):
            references.setdefault(slug_match, []).append(str(md_file.relative_to(root)))

    backlinks_dir = meta_dir / "backlinks"
    ensure_dirs(backlinks_dir)

    for slug, ref_list in references.items():
        bl_path = backlinks_dir / f"{slug}.json"
        bl_data = {"slug": slug, "backlinks": ref_list, "count": len(ref_list)}
        bl_path.write_text(json.dumps(bl_data, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
        if slug in docs_by_slug:
            docs_by_slug[slug]["referencedBy"] = ref_list

    save_catalog(catalog, catalog_path)
    print(f"Generated backlinks for {len(references)} documents.")
    print(f"Stored in _meta/backlinks/")

    for slug, ref_list in references.items():
        entry = docs_by_slug.get(slug)
        if not entry:
            continue
        doc_path = root / entry["path"]
        if not doc_path.exists():
            continue
        content = doc_path.read_text(encoding="utf-8")
        metadata, body = extract_frontmatter(content)
        if "<!-- backlinks -->" in body:
            continue
        ref_names = []
        for ref_path in ref_list[:10]:
            ref_slug = slug_from_path(ref_path)
            ref_entry = get_doc_entry(ref_slug, catalog)
            ref_title = ref_entry["title"] if ref_entry else ref_slug.replace("-", " ").title()
            ref_names.append(f"- [{ref_title}]({ref_path})")
        if ref_names:
            more = "" if len(ref_list) <= 10 else f"\n- ... and {len(ref_list) - 10} more"
            bl_section = f"\n<!-- backlinks -->\n\n## Referenced by\n\n{''.join(ref_names)}{more}\n"
            write_frontmatter(doc_path, metadata, body.rstrip() + bl_section)

    print("Appended 'Referenced by' sections.")
    return 0


def cmd_wiki_links_check(root: Path, meta_dir: Path, catalog_path: Path, fail_on_broken: bool) -> int:
    catalog = load_catalog(catalog_path)
    valid_slugs = {d["slug"] for d in catalog["docs"]}
    broken = []
    all_md = find_md_files(root)

    for md_file in all_md:
        content = md_file.read_text(encoding="utf-8")
        for slug_match, section_match in SLUG_RE.findall(content):
            if slug_match not in valid_slugs:
                broken.append((str(md_file.relative_to(root)), slug_match, section_match))

    if broken:
        print(f"Found {len(broken)} broken link(s):")
        for fp, slug, sec in broken:
            print(f"  {fp}: [[{slug}]]{'#' + sec if sec else ''}")
        return 1
    else:
        print(f"All links valid ({len(all_md)} files checked).")
        return 0


def cmd_wiki_migrate(root: Path, meta_dir: Path, catalog_path: Path, target: str) -> int:
    if target != "structured":
        print(f"Error: Unknown target: {target}", file=sys.stderr)
        return 1

    catalog = load_catalog(catalog_path)
    for cat, dirs in STRUCTURED_DIRS.items():
        for d in dirs:
            (root / cat / d).mkdir(parents=True, exist_ok=True)

    moved = unmoved = 0
    for doc in catalog["docs"]:
        current_path = root / doc["path"]
        if not current_path.exists():
            unmoved += 1
            continue
        status = doc.get("status", "draft")
        target_base = {"stable": root / "core", "deprecated": root / "archive"}.get(status, root / "drafts")
        rel = Path(doc["path"])
        target_path = target_base / rel
        if target_path != current_path:
            target_path.parent.mkdir(parents=True, exist_ok=True)
            shutil.move(str(current_path), str(target_path))
            doc["path"] = str(target_path.relative_to(root))
            moved += 1
        else:
            unmoved += 1

    save_catalog(catalog, catalog_path)
    print(f"Moved {moved}, unchanged {unmoved}. Structure: core/ reference/ drafts/ archive/")
    return 0


def cmd_wiki_status(root: Path, meta_dir: Path, catalog_path: Path) -> int:
    catalog = load_catalog(catalog_path)
    docs = catalog.get("docs", [])
    if not docs:
        print("No documents indexed. Run `repo-wiki.py wiki add --all` first.")
        return 0

    status_counts: dict[str, int] = {}
    category_counts: dict[str, int] = {}
    for doc in docs:
        s = doc.get("status", "unknown")
        c = doc.get("category", "unknown")
        status_counts[s] = status_counts.get(s, 0) + 1
        category_counts[c] = category_counts.get(c, 0) + 1

    print(f"Total documents: {len(docs)}\n")
    print("By status:")
    for s, c in sorted(status_counts.items()):
        print(f"  {s}: {c}")
    print("\nBy category:")
    for c, n in sorted(category_counts.items()):
        print(f"  {c}: {n}")

    needs_fm = sum(1 for f in find_md_files(root) if not read_md(f)[0])
    if needs_fm > 0:
        print(f"\n⚠ {needs_fm} file(s) without frontmatter.")
    return 0


# ===========================================================================
# Convert Commands (Binary → Text)
# ===========================================================================

def _convert_pdf(path: Path, out_path: Path) -> None:
    import fitz  # noqa: F811
    doc = fitz.open(str(path))
    lines = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        lines.append(f"## Page {page_num + 1}\n\n{text}\n")
    out_path.write_text("\n".join(lines), encoding="utf-8")
    doc.close()


def _convert_docx(path: Path, out_path: Path) -> None:
    import docx  # noqa: F811
    doc = docx.Document(str(path))
    lines = []
    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            level = int(para.style.name[-1]) if para.style.name[-1:].isdigit() else 1
            lines.append(f"{'#' * level} {para.text}\n")
        else:
            lines.append(f"{para.text}\n")
    for table in doc.tables:
        lines.append("\n| " + " | ".join(cell.text for cell in table.rows[0].cells) + " |\n")
        lines.append("| " + " | ".join("---" for _ in table.rows[0].cells) + " |\n")
        for row in table.rows[1:]:
            lines.append("| " + " | ".join(cell.text for cell in row.cells) + " |\n")
    out_path.write_text("\n".join(lines), encoding="utf-8")


def _convert_pptx(path: Path, out_path: Path) -> None:
    from pptx import Presentation  # noqa: F811
    prs = Presentation(str(path))
    lines = []
    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {slide_num}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    lines.append(f"  {para.text}\n")
    out_path.write_text("\n".join(lines), encoding="utf-8")


def _convert_xlsx(path: Path, out_path: Path) -> None:
    import openpyxl  # noqa: F811
    wb = openpyxl.load_workbook(str(path), data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## Sheet: {sheet_name}\n")
        rows = list(ws.iter_rows(values_only=True))
        if rows:
            lines.append("| " + " | ".join(str(c) for c in rows[0]) + " |\n")
            lines.append("| " + " | ".join("---" for _ in rows[0]) + " |\n")
            for row in rows[1:]:
                lines.append("| " + " | ".join(str(c or "") for c in row) + " |\n")
    out_path.write_text("\n".join(lines), encoding="utf-8")


def cmd_convert(root: Path, formats: list[str], out_dir: Optional[str], overwrite: bool) -> int:
    print(f"Converting binary documents ({', '.join(formats)})...")
    converter = {"pdf": _convert_pdf, "docx": _convert_docx, "pptx": _convert_pptx, "xlsx": _convert_xlsx}
    binary_files = find_binary_files(root, formats)
    if not binary_files:
        print("No binary files found.")
        return 0

    target_dir = Path(out_dir) if out_dir else root
    ensure_dirs(target_dir)
    converted = 0

    for bf in binary_files:
        rel = bf.relative_to(root)
        out_path = target_dir / f"{rel.stem}.md"
        if out_path.exists() and not overwrite:
            print(f"  SKIP {rel} (exists, use --overwrite)")
            continue
        try:
            converter[bf.suffix.lstrip(".")](bf, out_path)
            converted += 1
            print(f"  ✓ {rel} → {out_path.relative_to(root)}")
        except Exception as e:
            print(f"  ✗ {rel}: {e}")

    print(f"Converted {converted}/{len(binary_files)} files.")
    return 0


# ===========================================================================
# Index Commands
# ===========================================================================

def cmd_index(root: Path, out_path: Optional[str], dry_run: bool, ignore_dirs: list[str]) -> int:
    print("Building index.md...")
    all_md = find_md_files(root)
    binary_files = find_binary_files(root)

    # File inventory
    dir_files: dict[str, list[tuple[str, int]]] = {}
    for md in all_md:
        rel = md.relative_to(root)
        parent = str(rel.parent)
        dir_files.setdefault(parent, []).append((rel.name, md.stat().st_size // 1024))

    # Binary list
    binary_list = []
    for bf in binary_files:
        rel = bf.relative_to(root)
        converted = f"{rel.stem}.md"
        size = bf.stat().st_size // 1024
        binary_list.append((bf.name, str(rel), size, converted))

    # Keyword map
    keyword_map: dict[str, list[str]] = {}
    for md in all_md:
        content = md.read_text(encoding="utf-8")
        rel = str(md.relative_to(root))
        for m in HEADING_RE.finditer(content):
            heading = m.group(2).strip()
            kw = heading_to_slug(heading)
            if len(kw) > 3:
                keyword_map.setdefault(kw, []).append(rel)

    lines = ["# Repository Index", f"> Generated at {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')} / {len(all_md)} files\n",
             "## File Inventory\n"]
    for dir_name in sorted(dir_files.keys()):
        lines.append(f"### 📁 {dir_name}\n")
        for fname, size in sorted(dir_files[dir_name]):
            rel = f"{dir_name}/{fname}" if dir_name != "." else fname
            lines.append(f"- [{fname}]({rel}) ({size} KB)\n")

    lines.append("## Binary Documents\n")
    lines.append("| Filename | Path | Size | Converted |\n|----------|------|------|-----------|\n")
    for name, path, size, converted in binary_list:
        lines.append(f"| {name} | {path} | {size} KB | [{converted}]({converted}) |\n")

    lines.append("## Keyword Map\n")
    for kw in sorted(keyword_map.keys())[:100]:
        lines.append(f"### {kw}\n")
        for f in keyword_map[kw][:20]:
            lines.append(f"- [{f}]({f})\n")

    output_path = Path(out_path) if out_path else root / "index.md"
    output_str = "\n".join(lines)

    if dry_run:
        print(f"Would generate {len(lines)} lines at {output_path.relative_to(root)}")
        return 0

    output_path.write_text(output_str, encoding="utf-8")
    print(f"Generated {output_path.relative_to(root)} ({len(all_md)} files, {len(binary_list)} binaries, {len(keyword_map)} keywords)")
    return 0


# ===========================================================================
# Unified Runner
# ===========================================================================

def cmd_run(root: Path, skip_convert: bool, skip_index: bool, dry_run: bool) -> int:
    print("=" * 60)
    print("Repo Wiki — Full Workflow")
    print("=" * 60)
    t0 = time.time()

    # Phase 1: Wiki
    print("\n[1/3] Wiki initialization...")
    meta_dir = root / "_meta"
    catalog_path = meta_dir / "catalog.json"
    if not catalog_path.exists():
        cmd_wiki_init(root, meta_dir, meta_dir / "backlinks", catalog_path, root / "_index.md", root / ".repo-wikirc")

    print("[1/3] Adding documents...")
    cmd_wiki_add(root, meta_dir, catalog_path, True, None)

    print("[1/3] Building backlinks...")
    cmd_wiki_backlinks(root, meta_dir, catalog_path, None)

    print("[1/3] Checking links...")
    cmd_wiki_links_check(root, meta_dir, catalog_path, False)

    if not skip_convert:
        print("\n[2/3] Converting binary documents...")
        cmd_convert(root, list(BINARY_FORMATS.keys()), None, False)

    if not skip_index:
        print("\n[3/3] Building index.md...")
        cmd_index(root, None, False, [])

    elapsed = time.time() - t0
    print(f"\n{'=' * 60}")
    print(f"Done in {elapsed:.1f}s")
    print(f"{'=' * 60}")
    return 0


# ===========================================================================
# Dependency Installer
# ===========================================================================

def cmd_install_deps(fmts: list[str]) -> int:
    import importlib
    print("=== Checking repo-wiki dependencies ===\n")
    all_ok = True
    for fmt in fmts:
        pkg_map = {"pdf": ("fitz", "pymupdf"), "docx": ("docx", "python-docx"),
                   "pptx": ("pptx", "python-pptx"), "xlsx": ("openpyxl", "openpyxl")}
        if fmt not in pkg_map:
            continue
        import_name, pip_name = pkg_map[fmt]
        try:
            importlib.import_module(import_name)
            print(f"  [{fmt:4s}] {pip_name:15s}  ✅ Installed")
        except ImportError:
            print(f"  [{fmt:4s}] {pip_name:15s}  ❌ Missing")
            all_ok = False
            print(f"         → Running pip install {pip_name}...")
            result = subprocess.run([sys.executable, "-m", "pip", "install", pip_name], capture_output=True, text=True)
            if result.returncode == 0:
                print(f"         ✅ Installed")
            else:
                print(f"         ❌ Failed — run manually: pip install {pip_name}")
    print()
    if all_ok:
        print("✅ All dependencies met.")
    else:
        print("⚠️  Some dependencies missing.")
        return 1
    return 0


# ===========================================================================
# CLI Entry Point
# ===========================================================================

def main() -> int:
    parser = argparse.ArgumentParser(description="repo-wiki.py — Unified repository wiki", prog="repo-wiki.py")
    sub = parser.add_subparsers(dest="command")

    # wiki
    wiki_p = sub.add_parser("wiki")
    wiki_sub = wiki_p.add_subparsers(dest="subcommand")

    wiki_sub.add_parser("init", help="Initialize wiki structure")
    add_p = wiki_sub.add_parser("add")
    add_p.add_argument("--all", action="store_true")
    add_p.add_argument("--file", type=str)
    link_p = wiki_sub.add_parser("link")
    link_p.add_argument("from_file")
    link_p.add_argument("to_slug")
    link_p.add_argument("--section", type=str)
    bl_p = wiki_sub.add_parser("backlinks")
    bl_p.add_argument("--build", action="store_true")
    bl_p.add_argument("--file", type=str)
    check_p = wiki_sub.add_parser("links")
    check_p.add_argument("--check", action="store_true")
    check_p.add_argument("--fail-on-broken", action="store_true")
    mig_p = wiki_sub.add_parser("migrate")
    mig_p.add_argument("--to", type=str)
    wiki_sub.add_parser("status")

    # convert
    conv_p = sub.add_parser("convert")
    conv_p.add_argument("--root", type=str, default=str(DEFAULT_DOCS_DIR))
    conv_p.add_argument("--formats", nargs="+", default=list(BINARY_FORMATS.keys()))
    conv_p.add_argument("--out", type=str)
    conv_p.add_argument("--overwrite", action="store_true")

    # index
    idx_p = sub.add_parser("index")
    idx_p.add_argument("--root", type=str, default=str(DEFAULT_DOCS_DIR))
    idx_p.add_argument("--out", type=str)
    idx_p.add_argument("--dry-run", action="store_true")
    idx_p.add_argument("--ignore-dirs", nargs="+", default=list(IGNORE_DIRS))

    # run
    run_p = sub.add_parser("run")
    run_p.add_argument("--root", type=str, default=str(DEFAULT_DOCS_DIR))
    run_p.add_argument("--skip-convert", action="store_true")
    run_p.add_argument("--skip-index", action="store_true")
    run_p.add_argument("--dry-run", action="store_true")

    # install-deps
    dep_p = sub.add_parser("install-deps")
    dep_p.add_argument("--fmt", nargs="+", default=list(BINARY_FORMATS.keys()))

    args = parser.parse_args()
    if not args.command:
        parser.print_help()
        return 1

    root = Path(args.root) if hasattr(args, "root") and args.root else DEFAULT_DOCS_DIR
    meta_dir = root / "_meta"
    catalog_path = meta_dir / "catalog.json"

    cmds = {
        ("wiki", "init"): lambda: cmd_wiki_init(root, meta_dir, meta_dir / "backlinks", catalog_path, root / "_index.md", root / ".repo-wikirc"),
        ("wiki", "add"): lambda: cmd_wiki_add(root, meta_dir, catalog_path, args.all, args.file),
        ("wiki", "link"): lambda: cmd_wiki_link(root, args.from_file, args.to_slug, args.section),
        ("wiki", "backlinks"): lambda: cmd_wiki_backlinks(root, meta_dir, catalog_path, args.file),
        ("wiki", "links"): lambda: cmd_wiki_links_check(root, meta_dir, catalog_path, args.fail_on_broken),
        ("wiki", "migrate"): lambda: cmd_wiki_migrate(root, meta_dir, catalog_path, args.to),
        ("wiki", "status"): lambda: cmd_wiki_status(root, meta_dir, catalog_path),
        ("convert", None): lambda: cmd_convert(root, args.formats, args.out, args.overwrite),
        ("index", None): lambda: cmd_index(root, args.out, args.dry_run, args.ignore_dirs),
        ("run", None): lambda: cmd_run(root, args.skip_convert, args.skip_index, args.dry_run),
        ("install-deps", None): lambda: cmd_install_deps(args.fmt),
    }

    key = (args.command, getattr(args, "subcommand", None))
    fn = cmds.get(key)
    if fn:
        return fn()
    print(f"Unknown command: {args.command}", file=sys.stderr)
    return 1


if __name__ == "__main__":
    sys.exit(main())
