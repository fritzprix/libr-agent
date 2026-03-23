from pptx import Presentation
import sys

def pptx_to_md(pptx_path, md_path):
    try:
        prs = Presentation(pptx_path)
        with open(md_path, 'w', encoding='utf-8') as f:
            for i, slide in enumerate(prs.slides):
                f.write(f"## Slide {i + 1}\n\n")
                
                # 슬라이드 제목 추출
                if slide.shapes.title and slide.shapes.title.text:
                    f.write(f"### {slide.shapes.title.text.strip()}\n\n")
                
                # 본문 내용 추출
                for shape in slide.shapes:
                    if shape == slide.shapes.title: 
                        continue
                    if hasattr(shape, "text") and shape.text.strip():
                        f.write(f"{shape.text.strip()}\n\n")
                f.write("---\n\n")
        print(f"✅ Successfully converted PPTX: {pptx_path} -> {md_path}")
    except Exception as e:
        print(f"❌ Error: {e}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python pptx_to_md.py <input.pptx> <output.md>")
    else:
        pptx_to_md(sys.argv[1], sys.argv[2])
