#!/usr/bin/env python3
"""
Workspace Binary Document Converter
====================================
Convert binary document files (PPTX, PDF, DOCX, XLSX) into text/Markdown.

Usage:
    python convert_binary_docs.py [--root ROOT_DIR] [--out OUT_DIR] [--formats pdf pptx docx xlsx]

Options:
    --root      Root directory to scan (default: current directory)
    --out       Output directory for converted files (default: next to the source file)
    --formats   File formats to convert (default: pdf pptx docx xlsx)
    --overwrite Overwrite existing converted files (default: False)
    --dry-run   Print target files without converting anything

Dependencies:
    pip install pymupdf python-docx python-pptx openpyxl

Examples:
    # Convert the full workspace
    python convert_binary_docs.py --root . --out ./converted

    # Convert PDFs only
    python convert_binary_docs.py --root . --formats pdf

    # Preview target files with dry-run
    python convert_binary_docs.py --root . --dry-run
"""

import argparse
import io
import sys
from pathlib import Path

# Windows에서 cp949 인코딩으로 인한 UnicodeEncodeError 방지
if sys.platform.startswith("win"):
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

# Cwd 독립적인 임포트를 위해 sys.path에 scripts 디렉토리 추가
sys.path.append(str(Path(__file__).resolve().parent))
from utils import get_source_from_md, get_path_hash

SKIP_DIRS = {".git", "__pycache__", "node_modules", ".github", "venv", ".venv", ".libragent", "attachments"}


def find_binary_files(root: Path, formats: list[str]) -> list[Path]:
    """Recursively find binary files matching the requested formats."""
    exts = {f".{fmt.lstrip('.')}" for fmt in formats}
    result = []
    for p in root.rglob("*"):
        if any(part in SKIP_DIRS for part in p.parts):
            continue
        if p.suffix.lower() in exts and p.is_file():
            result.append(p)
    return sorted(result)


def convert_pdf(src: Path) -> str:
    """Convert PDF content to Markdown text."""
    try:
        import fitz  # pymupdf
    except ImportError:
        return f"[ERROR] pymupdf not installed. Run: pip install pymupdf\nSource: {src}"

    doc = fitz.open(str(src))
    pages = []
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text:
            pages.append(f"## Page {i}\n\n{text}")
    doc.close()
    return "\n\n---\n\n".join(pages) if pages else "(no text content)"


def convert_docx(src: Path) -> str:
    """Convert DOCX content to Markdown text."""
    try:
        from docx import Document
    except ImportError:
        return f"[ERROR] python-docx not installed. Run: pip install python-docx\nSource: {src}"

    doc = Document(str(src))
    lines = []
    for para in doc.paragraphs:
        style = para.style.name.lower()
        text = para.text.strip()
        if not text:
            lines.append("")
            continue
        if style.startswith("heading 1"):
            lines.append(f"# {text}")
        elif style.startswith("heading 2"):
            lines.append(f"## {text}")
        elif style.startswith("heading 3"):
            lines.append(f"### {text}")
        else:
            lines.append(text)

    # Tables
    for table in doc.tables:
        if not table.rows:
            continue
        header = [cell.text.strip() for cell in table.rows[0].cells]
        lines.append("\n| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in table.rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")

    return "\n".join(lines)


def convert_pptx(src: Path) -> str:
    """Convert PPTX content to Markdown text."""
    try:
        from pptx import Presentation
    except ImportError:
        return f"[ERROR] python-pptx not installed. Run: pip install python-pptx\nSource: {src}"

    prs = Presentation(str(src))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
    return "\n\n---\n\n".join(slides) if slides else "(no text content)"


def convert_xlsx(src: Path) -> str:
    """Convert XLSX sheets into Markdown tables."""
    try:
        import openpyxl
    except ImportError:
        return f"[ERROR] openpyxl not installed. Run: pip install openpyxl\nSource: {src}"

    wb = openpyxl.load_workbook(str(src), read_only=True, data_only=True)
    sections = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        # Skip empty sheets
        non_empty = [r for r in rows if any(c is not None for c in r)]
        if not non_empty:
            continue
        lines = [f"## Sheet: {sheet_name}\n"]
        for idx, row in enumerate(non_empty):
            cells = [str(c) if c is not None else "" for c in row]
            lines.append("| " + " | ".join(cells) + " |")
            if idx == 0:
                lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
        sections.append("\n".join(lines))
    wb.close()
    return "\n\n---\n\n".join(sections) if sections else "(no content)"


CONVERTERS = {
    ".pdf": convert_pdf,
    ".docx": convert_docx,
    ".pptx": convert_pptx,
    ".xlsx": convert_xlsx,
}


def convert_file(src: Path, out_dir: Path | None, overwrite: bool, used_dests: set[str]) -> tuple[Path, str]:
    """Convert a single file and return (output_path, status_message)."""
    ext = src.suffix.lower()
    converter = CONVERTERS.get(ext)
    if not converter:
        return src, "SKIP (no converter)"

    dest_dir = out_dir if out_dir else src.parent
    
    base_dest = dest_dir / (src.stem + ".md")
    dest = base_dest
    dest_key = str(dest.resolve().as_posix())
    
    # Resolve unique filename by appending path hash if conflict occurs
    if dest_key in used_dests or (dest.exists() and not overwrite):
        source_in_md = get_source_from_md(dest)
        if source_in_md and Path(source_in_md).resolve() == src.resolve():
            # Same source, safe to skip
            used_dests.add(dest_key)
            return dest, "SKIP (exists)"
        else:
            # Different source (collision) -> append short hash of source path
            path_hash = get_path_hash(src)
            dest = dest_dir / f"{src.stem}_{path_hash}.md"
            dest_key = str(dest.resolve().as_posix())
            
            # Resolve potential hash collisions
            counter = 1
            while True:
                if dest_key in used_dests:
                    dest = dest_dir / f"{src.stem}_{path_hash}_{counter}.md"
                    dest_key = str(dest.resolve().as_posix())
                    counter += 1
                    continue
                if dest.exists() and not overwrite:
                    source_in_md = get_source_from_md(dest)
                    if source_in_md and Path(source_in_md).resolve() == src.resolve():
                        used_dests.add(dest_key)
                        return dest, "SKIP (exists)"
                    else:
                        dest = dest_dir / f"{src.stem}_{path_hash}_{counter}.md"
                        dest_key = str(dest.resolve().as_posix())
                        counter += 1
                        continue
                break

    dest_dir.mkdir(parents=True, exist_ok=True)

    # Performance logging for large files (e.g. pptx > 5MB)
    size_mb = src.stat().st_size / (1024 * 1024)
    if size_mb > 5:
        print(f"       ⌛ Converting large file ({size_mb:.1f} MB)... this may take a moment.")

    content = converter(src)
    header = f"# {src.name}\n\n> Source: `{src}`  \n> Converted at: {_now()}\n\n---\n\n"
    dest.write_text(header + content, encoding="utf-8")
    used_dests.add(dest_key)
    return dest, "OK"


def _now() -> str:
    from datetime import datetime
    return datetime.now().strftime("%Y-%m-%d %H:%M")


def main():
    parser = argparse.ArgumentParser(description="Binary document to Markdown converter")
    parser.add_argument("--root", default=".", help="Root directory to scan")
    parser.add_argument("--out", default=None, help="Output directory (defaults to the source file location)")
    parser.add_argument("--formats", nargs="+", default=["pdf", "pptx", "docx", "xlsx"])
    parser.add_argument("--overwrite", action="store_true")
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args()

    root = Path(args.root).resolve()
    out_dir = Path(args.out).resolve() if args.out else None

    files = find_binary_files(root, args.formats)
    print(f"Found files: {len(files)}\n")

    if args.dry_run:
        for f in files:
            print(f"  {f.relative_to(root)}")
        return

    results = {"OK": 0, "SKIP (exists)": 0, "SKIP (no converter)": 0, "ERROR": 0}
    used_dests = set()
    for src in files:
        try:
            dest, status = convert_file(src, out_dir, args.overwrite, used_dests)
            results[status] = results.get(status, 0) + 1
            icon = "✅" if status == "OK" else "⏭️"
            print(f"  {icon} [{status}] {src.relative_to(root)}")
            if status == "OK":
                print(f"       → {dest}")
        except Exception as e:
            results["ERROR"] += 1
            print(f"  ❌ [ERROR] {src.relative_to(root)}: {e}")

    print(f"\nDone: {results}")


if __name__ == "__main__":
    main()
